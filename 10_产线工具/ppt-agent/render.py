"""
PPT Agent Unified Renderer
Usage: python render.py <style_id> <content_md> <output_pptx>
Example: python render.py b-dark-cyber workfiles/PPT-xxx_content.md workfiles/PPT-xxx.pptx
"""

import yaml, re, os, sys, random, math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = os.path.dirname(os.path.abspath(__file__))

def rgb(h):
    return RGBColor(int(h[1:3],16), int(h[3:5],16), int(h[5:7],16))

def hex_to_rgb(hex_str):
    """From hex string like '#ff0000' to (255,0,0)"""
    h = hex_str.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0,2,4))

# ============================================================
# STYLE-SPECIFIC DECORATION FUNCTIONS
# ============================================================

def add_hex_grid(slide, prs, PAL):
    """Hardcore: hexagonal grid overlay"""
    sz = Pt(30)
    for r in range(int(prs.slide_height / (sz*1.6)) + 2):
        ox = (sz * 1.8 / 2) if r % 2 == 1 else 0
        for c in range(int(prs.slide_width / (sz*1.8)) + 2):
            h = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, c*sz*1.8+ox, r*sz*1.6, sz, sz)
            h.fill.solid(); h.fill.fore_color.rgb = rgb(PAL["background"])
            h.line.color.rgb = rgb(PAL["divider"]); h.line.width = Pt(0.3)

def add_rivet_divider(slide, left, top, width, PAL):
    """Hardcore: rivet line divider"""
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    ln.fill.solid(); ln.fill.fore_color.rgb = rgb(PAL["divider"]); ln.line.fill.background()
    for i in range(1, 7):
        x, y = left + width/7*i - Pt(5), top - Pt(3)
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Pt(10), Pt(10))
        d.fill.solid(); d.fill.fore_color.rgb = rgb(PAL["primary"]); d.line.fill.background()
        id = slide.shapes.add_shape(MSO_SHAPE.OVAL, x+Pt(2), y+Pt(2), Pt(6), Pt(6))
        id.fill.solid(); id.fill.fore_color.rgb = rgb(PAL.get("warning_stripe_fg","#f5a623"))
        id.line.fill.background()

def add_warning_stripes(slide, left, top, width, height, PAL):
    """Hardcore: diagonal warning stripe bar"""
    sh, gap = Pt(6), Pt(4)
    step = sh + gap
    for i in range(int(height/step)+4):
        s = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, left+i*step-height, top, sh, height*2)
        s.fill.solid()
        s.fill.fore_color.rgb = rgb(PAL["warning_stripe_bg"]) if i%2==0 else rgb(PAL["warning_stripe_fg"])
        s.line.fill.background(); s.rotation = -45.0

def add_metallic_title(slide, left, top, width, height, text, PAL, TYPO, fs=38):
    """Hardcore: metallic effect title"""
    hf = TYPO["heading"].split(",")[0].strip()
    sh = slide.shapes.add_textbox(left+Pt(2), top+Pt(2), width, height)
    sh.text_frame.word_wrap = True
    sh.text_frame.paragraphs[0].text = text.upper(); sh.text_frame.paragraphs[0].font.size = Pt(fs)
    sh.text_frame.paragraphs[0].font.color.rgb = rgb(PAL["divider"])
    sh.text_frame.paragraphs[0].font.bold = True; sh.text_frame.paragraphs[0].font.name = hf

    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True; mp = tb.text_frame.paragraphs[0]
    mp.text = text.upper(); mp.font.size = Pt(fs); mp.font.color.rgb = rgb(PAL["primary"])
    mp.font.bold = True; mp.font.name = hf
    return tb

def add_hardcore_section_header(slide, left, top, width, text, PAL, TYPO):
    """Hardcore: red bracket section header"""
    br = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(4), Inches(0.35))
    br.fill.solid(); br.fill.fore_color.rgb = rgb(PAL["accent"]); br.line.fill.background()
    tb = slide.shapes.add_textbox(left+Pt(12), top-Pt(2), width-Pt(24), Inches(0.4))
    tb.text_frame.word_wrap = True; pp = tb.text_frame.paragraphs[0]
    pp.text = text; pp.font.size = Pt(22); pp.font.color.rgb = rgb(PAL["warning_stripe_fg"])
    pp.font.bold = True; pp.font.name = TYPO["body"].split(",")[0].strip()

def add_diagonal_stripes_bg(slide, prs, colors, stripe_w=Inches(2.0)):
    """Memphis: bright diagonal stripe background"""
    n = int(prs.slide_width / stripe_w) + 10
    for i in range(n):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, i*stripe_w-prs.slide_height, 0, stripe_w*1.5, prs.slide_height*2)
        s.fill.solid(); s.fill.fore_color.rgb = rgb(colors[i%len(colors)])
        s.line.fill.background(); s.rotation = -35.0

def add_terrazzo_speckle(slide, prs, PAL, count=60):
    """Memphis: terrazzo speckle overlay"""
    sc = [PAL["terrazzo_speckle_1"], PAL["terrazzo_speckle_2"], PAL["terrazzo_speckle_3"], PAL["terrazzo_speckle_4"]]
    for _ in range(count):
        d = slide.shapes.add_shape(random.choice([MSO_SHAPE.OVAL, MSO_SHAPE.RECTANGLE, MSO_SHAPE.DIAMOND]),
                                    random.randint(0, int(prs.slide_width)),
                                    random.randint(0, int(prs.slide_height)),
                                    Pt(random.randint(3,12)), Pt(random.randint(3,12)))
        d.fill.solid(); d.fill.fore_color.rgb = rgb(random.choice(sc)); d.line.fill.background()
        d.rotation = random.uniform(-30, 30)

def add_scattered_accents(slide, prs, colors, count=5):
    """Memphis: scattered decorative shapes"""
    shapes_pool = [MSO_SHAPE.OVAL, MSO_SHAPE.DIAMOND, MSO_SHAPE.ISOSCELES_TRIANGLE, MSO_SHAPE.STAR_5_POINT]
    for _ in range(count):
        s = slide.shapes.add_shape(random.choice(shapes_pool),
                                    random.randint(int(Inches(0.5)), int(prs.slide_width-Inches(1.5))),
                                    random.randint(int(Inches(0.5)), int(prs.slide_height-Inches(1.5))),
                                    Pt(random.randint(15,40)), Pt(random.randint(15,40)))
        s.fill.solid(); s.fill.fore_color.rgb = rgb(random.choice(colors))
        s.line.fill.background(); s.rotation = random.uniform(-45, 45)

def add_outline_title(slide, left, top, width, height, text, fill_c, stroke_c, TYPO, fs=38):
    """Memphis: outline-stroke title"""
    hf = TYPO["heading"].split(",")[0].strip()
    for ox,oy in [(3,3),(-3,3),(3,-3),(-3,-3),(3,0),(-3,0),(0,3),(0,-3)]:
        sh = slide.shapes.add_textbox(left+Pt(ox), top+Pt(oy), width, height)
        sh.text_frame.word_wrap = True; sp = sh.text_frame.paragraphs[0]
        sp.text = text; sp.font.size = Pt(fs); sp.font.color.rgb = rgb(stroke_c)
        sp.font.bold = True; sp.font.name = hf; sp.alignment = PP_ALIGN.LEFT
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True; mp = tb.text_frame.paragraphs[0]
    mp.text = text; mp.font.size = Pt(fs); mp.font.color.rgb = rgb(fill_c)
    mp.font.bold = True; mp.font.name = hf; mp.alignment = PP_ALIGN.LEFT
    return tb

def add_wavy_divider(slide, left, top, width, colors):
    """Memphis: wavy dotted divider"""
    for i in range(int(width/(Pt(16)))):
        x = left + i*Pt(16); y = top + math.sin(i*0.8)*Pt(6)+Pt(4)
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Pt(6), Pt(6))
        d.fill.solid(); d.fill.fore_color.rgb = rgb(colors[i%len(colors)])
        d.line.fill.background()

def set_round_corners(shape, radius_emu):
    """Apply rounded corners via XML"""
    prstGeom = shape._element.find(qn('a:prstGeom'))
    if prstGeom is None: return
    avLst = prstGeom.find(qn('a:avLst'))
    if avLst is not None: prstGeom.remove(avLst)
    prstGeom.set('prst', 'roundRect')
    avLst = prstGeom.makeelement(qn('a:avLst'), {})
    gd = avLst.makeelement(qn('a:gd'), {'name':'adj', 'fmla': f'val {radius_emu}'})
    avLst.append(gd); prstGeom.append(avLst)

def add_double_border_card(slide, left, top, width, height, PAL):
    """Memphis: double-line rounded card"""
    o = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    o.fill.solid(); o.fill.fore_color.rgb = rgb(PAL["surface"])
    o.line.color.rgb = rgb(PAL["primary"]); o.line.width = Pt(3)
    set_round_corners(o, Emu(Inches(0.2)))
    i = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left+Pt(6), top+Pt(6), width-Pt(12), height-Pt(12))
    i.fill.solid(); i.fill.fore_color.rgb = rgb(PAL["surface"])
    i.line.color.rgb = rgb(PAL["secondary_1"]); i.line.width = Pt(1.5)
    set_round_corners(i, Emu(Inches(0.18)))
    return (left+Pt(8), top+Pt(8), width-Pt(16), height-Pt(16))

# ============================================================
# COMMON HELPERS
# ============================================================

def parse_content(content_file):
    """Parse content.md into list of slide dicts"""
    with open(content_file, encoding="utf-8") as f:
        txt = f.read()
    blocks = []
    for b in re.split(r"\n---\n", txt):
        m = re.match(r"# 第(\d+)页\s*[·●•]\s*(.+)", b.strip())
        if m:
            body = b[b.index("\n"):].strip()
            body = re.sub(r"^【.*?】\n", "", body, flags=re.MULTILINE)
            blocks.append({"n": int(m.group(1)), "t": m.group(2), "b": body})
    return blocks

def add_simple_title(slide, left, top, width, height, text, PAL, TYPO, fs=36, align=PP_ALIGN.LEFT):
    """Generic title for simple styles"""
    hf = TYPO["heading"].split(",")[0].strip()
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True; mp = tb.text_frame.paragraphs[0]
    mp.text = text; mp.font.size = Pt(fs); mp.font.color.rgb = rgb(PAL["text_primary"])
    mp.font.bold = True; mp.font.name = hf; mp.alignment = align
    return tb

def add_body_text(slide, left, top, width, body_text, PAL, TYPO):
    """Parse markdown body text and render as textboxes"""
    bfont = TYPO["body"].split(",")[0].strip()
    cur = top
    secs = re.split(r"\n(?=\*\*)", body_text)
    for sec in secs:
        sec = sec.strip()
        if not sec: continue
        lines = sec.split("\n")
        stitle = None
        if lines[0].startswith("**") and "**" in lines[0][2:]:
            stitle = lines[0].strip("* "); lines = lines[1:]
        clean = [l.strip() for l in lines if l.strip() and not l.startswith(">") and not l.startswith("|:---")]
        if not clean and not stitle: continue

        if stitle:
            ss = slide.shapes.add_textbox(left, cur, width, Inches(0.4))
            ss.text_frame.word_wrap = True; sp = ss.text_frame.paragraphs[0]
            sp.text = stitle; sp.font.size = Pt(20); sp.font.color.rgb = rgb(PAL["primary"])
            sp.font.bold = True; sp.font.name = bfont
            cur += Inches(0.45)

        cl = [l.lstrip("- ") for l in clean]
        if cl:
            th = Pt(20) * len(cl) + Inches(0.1)
            bx = slide.shapes.add_textbox(left+Inches(0.2), cur, width-Inches(0.2), th)
            bx.text_frame.word_wrap = True
            for i, line in enumerate(cl):
                pp = bx.text_frame.paragraphs[0] if i==0 else bx.text_frame.add_paragraph()
                pp.text = line; pp.font.size = Pt(14); pp.font.color.rgb = rgb(PAL["text_primary"])
                pp.font.name = bfont; pp.space_after = Pt(5)
            cur += th + Inches(0.3)
    return cur

def render_slide_common(prs, sd, PAL, TYPO, style_id, is_hardcore=False, is_memphis=False):
    """Common slide building logic"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(PAL["background"])
    return slide

# ============================================================
# STYLE DISPATCH
# ============================================================

def render_standard_style(prs, blocks, PAL, TYPO):
    """For styles a, c, d, e, f: clean professional look"""
    MG, BY = Inches(0.8), Inches(1.6)
    BW = prs.slide_width - MG*2
    bfont = TYPO["body"].split(",")[0].strip()

    for sd in blocks:
        slide = render_slide_common(prs, sd, PAL, TYPO, "standard")

        # Page number
        sn = slide.shapes.add_textbox(MG, Inches(0.1), Inches(1.5), Inches(0.4))
        sn.text_frame.paragraphs[0].text = f"0{sd['n']}" if sd['n']<10 else str(sd['n'])
        sn.text_frame.paragraphs[0].font.size = Pt(22)
        sn.text_frame.paragraphs[0].font.color.rgb = rgb(PAL["primary"])
        sn.text_frame.paragraphs[0].font.bold = True

        # Title
        add_simple_title(slide, MG, Inches(0.4), BW, Inches(0.8), sd["t"], PAL, TYPO, fs=36)

        # Body
        add_body_text(slide, MG, BY, BW, sd["b"], PAL, TYPO)

        # Bottom accent line
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MG, prs.slide_height-Inches(0.3), Inches(1.0), Pt(3))
        ln.fill.solid(); ln.fill.fore_color.rgb = rgb(PAL["primary"]); ln.line.fill.background()

def render_f_style(prs, blocks, PAL, TYPO):
    """Style f: elegant classic with gold accent"""
    MG, BY = Inches(1.0), Inches(1.8)
    BW = prs.slide_width - MG*2

    for sd in blocks:
        slide = render_slide_common(prs, sd, PAL, TYPO, "standard")

        # Top thin gold line
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MG, Inches(0.2), prs.slide_width-MG*2, Pt(1.5))
        ln.fill.solid(); ln.fill.fore_color.rgb = rgb(PAL["accent"]); ln.line.fill.background()

        sn = slide.shapes.add_textbox(MG, Inches(0.35), Inches(1.5), Inches(0.4))
        sn.text_frame.paragraphs[0].text = f"0{sd['n']}" if sd['n']<10 else str(sd['n'])
        sn.text_frame.paragraphs[0].font.size = Pt(20)
        sn.text_frame.paragraphs[0].font.color.rgb = rgb(PAL["accent"])
        sn.text_frame.paragraphs[0].font.bold = True

        add_simple_title(slide, MG, Inches(0.65), BW, Inches(0.9), sd["t"], PAL, TYPO, fs=36)
        add_body_text(slide, MG, BY, BW, sd["b"], PAL, TYPO)

        # Bottom gold line
        ln2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MG, prs.slide_height-Inches(0.45), prs.slide_width-MG*2, Pt(1.5))
        ln2.fill.solid(); ln2.fill.fore_color.rgb = rgb(PAL["accent"]); ln2.line.fill.background()

def render_b_style(prs, blocks, PAL, TYPO):
    """Style b: dark cyber with glass cards and glow"""
    MG, BY = Inches(0.8), Inches(1.6)
    BW = prs.slide_width - MG*2

    for sd in blocks:
        slide = render_slide_common(prs, sd, PAL, TYPO, "standard")

        # Corner accent
        if TYPO.get("layout", {}).get("corner_accent"):
            s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, prs.slide_width-Inches(1.5), prs.slide_height-Inches(1.5), Inches(1.5), Inches(1.5))
            s.fill.solid(); s.fill.fore_color.rgb = rgb(PAL["primary"]); s.line.fill.background()

        tn = f"0{sd['n']}" if sd['n']<10 else str(sd['n'])
        sn = slide.shapes.add_textbox(MG, Inches(0.1), Inches(1.5), Inches(0.4))
        sn.text_frame.paragraphs[0].text = tn
        sn.text_frame.paragraphs[0].font.size = Pt(24)
        sn.text_frame.paragraphs[0].font.color.rgb = rgb(PAL["primary"])
        sn.text_frame.paragraphs[0].font.bold = True

        add_simple_title(slide, MG, Inches(0.4), BW, Inches(0.8), sd["t"], PAL, TYPO, fs=36)

        # Glass card for body
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MG, BY, BW, prs.slide_height-BY-Inches(0.8))
        card.fill.solid(); card.fill.fore_color.rgb = rgb(PAL["surface"])
        card.line.color.rgb = rgb(PAL["divider"]); card.line.width = Pt(1)
        add_body_text(slide, MG+Inches(0.2), BY+Inches(0.1), BW-Inches(0.4), sd["b"], PAL, TYPO)

def render_g_style(prs, blocks, PAL, TYPO):
    """Style g: hardcore industrial"""
    MG, BY = Inches(0.7), Inches(2.0)
    BW = prs.slide_width - MG*2
    hfont = TYPO["heading"].split(",")[0].strip()

    for sd in blocks:
        slide = render_slide_common(prs, sd, PAL, TYPO, "standard")
        add_hex_grid(slide, prs, PAL)
        add_warning_stripes(slide, 0, Inches(0.15), prs.slide_width, Inches(1.3), PAL)
        add_metallic_title(slide, MG, Inches(0.35), BW-MG, Inches(0.8), sd["t"], PAL, TYPO, fs=38)

        # Badge
        bx = prs.slide_width - MG - Inches(0.8)
        badge = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, bx, Inches(0.45), Inches(0.6), Inches(0.6))
        badge.fill.solid(); badge.fill.fore_color.rgb = rgb(PAL["accent"]); badge.line.fill.background()
        bp = badge.text_frame.paragraphs[0]
        bp.text = f"0{sd['n']}" if sd['n']<10 else str(sd['n'])
        bp.font.size = Pt(20); bp.font.color.rgb = RGBColor(255,255,255)
        bp.font.bold = True; bp.font.name = hfont; bp.alignment = PP_ALIGN.CENTER

        add_rivet_divider(slide, MG, Inches(1.55), BW, PAL)

        # Body with hardcore section headers
        cur = BY; bfont = TYPO["body"].split(",")[0].strip()
        secs = re.split(r"\n(?=\*\*)", sd["b"])
        for sec in secs:
            sec = sec.strip()
            if not sec: continue
            lines = sec.split("\n")
            stitle = None
            if lines[0].startswith("**") and "**" in lines[0][2:]:
                stitle = lines[0].strip("* "); lines = lines[1:]
            clean = [l.strip() for l in lines if l.strip() and not l.startswith(">") and not l.startswith("|:---")]
            if not clean and not stitle: continue

            if stitle:
                add_hardcore_section_header(slide, MG+Inches(0.15), cur, BW-Inches(0.3), stitle, PAL, TYPO)
                cur += Inches(0.55)

            cl = [l.lstrip("- ") for l in clean]
            if cl:
                th = Pt(20) * len(cl) + Inches(0.15)
                bx = slide.shapes.add_textbox(MG+Inches(0.4), cur, BW-Inches(0.8), th)
                bx.text_frame.word_wrap = True
                for i, line in enumerate(cl):
                    pp = bx.text_frame.paragraphs[0] if i==0 else bx.text_frame.add_paragraph()
                    pp.text = line; pp.font.size = Pt(14)
                    pp.font.color.rgb = rgb(PAL["text_primary"])
                    pp.font.name = bfont; pp.space_after = Pt(5)
                cur += th + Inches(0.35)

        add_rivet_divider(slide, MG, prs.slide_height-Inches(0.25), BW, PAL)

def render_h_style(prs, blocks, PAL, TYPO):
    """Style h: memphis dopamine"""
    MG, BY = Inches(0.8), Inches(2.0)
    BW = prs.slide_width - MG*2
    ACCS = [PAL["primary"], PAL["accent"], PAL["secondary_1"], PAL["secondary_2"],
            PAL["secondary_3"], PAL["secondary_4"]]
    hfont = TYPO["heading"].split(",")[0].strip()
    stripe_sets = [
        [PAL["primary"], PAL["secondary_1"], PAL["accent"], PAL["secondary_3"]],
        [PAL["secondary_2"], PAL["accent"], PAL["secondary_4"], PAL["secondary_1"]],
        [PAL["primary"], PAL["secondary_2"], PAL["secondary_3"], PAL["accent"]],
    ]

    for idx, sd in enumerate(blocks):
        slide = render_slide_common(prs, sd, PAL, TYPO, "standard")
        add_diagonal_stripes_bg(slide, prs, stripe_sets[idx%3])
        add_terrazzo_speckle(slide, prs, PAL, 60)
        add_scattered_accents(slide, prs, ACCS, 5)

        # Outline title with random offset
        tx = MG + Pt(random.randint(-10, 15))
        ty = Inches(0.4) + Pt(random.randint(-5, 10))
        tw = BW - Pt(random.randint(0, 30))
        fill_c = ACCS[idx % len(ACCS)]
        add_outline_title(slide, tx, ty, tw, Inches(1.0), sd["t"], fill_c, PAL["stroke_outline"], TYPO, fs=38)

        # Bubble page number
        bubble = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         prs.slide_width-MG-Inches(1.3), Inches(0.5),
                                         Inches(1.0), Inches(0.55))
        bubble.fill.solid(); bubble.fill.fore_color.rgb = rgb(PAL["primary"])
        bubble.line.color.rgb = rgb(PAL["stroke_outline"]); bubble.line.width = Pt(2.5)
        set_round_corners(bubble, Emu(Inches(0.3)))
        bp = bubble.text_frame.paragraphs[0]
        bp.text = f"0{sd['n']}" if sd['n']<10 else str(sd['n'])
        bp.font.size = Pt(22); bp.font.color.rgb = RGBColor(255,255,255)
        bp.font.bold = True; bp.font.name = hfont; bp.alignment = PP_ALIGN.CENTER

        add_wavy_divider(slide, MG, Inches(1.55), BW, ACCS)

        # Body with memphis decoration
        cur = BY; bfont = TYPO["body"].split(",")[0].strip()
        secs = re.split(r"\n(?=\*\*)", sd["b"])
        for sec in secs:
            sec = sec.strip()
            if not sec: continue
            lines = sec.split("\n")
            stitle = None
            if lines[0].startswith("**") and "**" in lines[0][2:]:
                stitle = lines[0].strip("* "); lines = lines[1:]
            clean = [l.strip() for l in lines if l.strip() and not l.startswith(">") and not l.startswith("|:---")]
            if not clean and not stitle: continue

            if stitle:
                sh_c = ACCS[(idx+1)%len(ACCS)]
                ss = slide.shapes.add_textbox(MG, cur, BW, Inches(0.45))
                ss.text_frame.word_wrap = True; sp = ss.text_frame.paragraphs[0]
                sp.text = f"✦ {stitle}"; sp.font.size = Pt(22)
                sp.font.color.rgb = rgb(sh_c); sp.font.bold = True; sp.font.name = hfont
                cur += Inches(0.45)
                # Dashed underline
                for x in range(int(MG), int(MG+BW), int(Pt(12))):
                    d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, cur-Pt(2), Pt(8), Pt(2))
                    d.fill.solid(); d.fill.fore_color.rgb = rgb(sh_c); d.line.fill.background()
                cur += Inches(0.15)

            cl = [l.lstrip("- ") for l in clean]
            if cl:
                card_h = Pt(21)*len(cl)+Inches(0.3)
                card_l, card_t, card_w, card_h2 = add_double_border_card(slide, MG, cur, BW, card_h, PAL)
                bx = slide.shapes.add_textbox(card_l+Inches(0.1), card_t, card_w, card_h2)
                bx.text_frame.word_wrap = True
                for i, line in enumerate(cl):
                    pp = bx.text_frame.paragraphs[0] if i==0 else bx.text_frame.add_paragraph()
                    pp.text = line; pp.font.size = Pt(14)
                    pp.font.color.rgb = rgb(PAL["text_primary"])
                    pp.font.name = bfont; pp.space_after = Pt(5)
                cur += card_h + Inches(0.25)

        add_wavy_divider(slide, MG, prs.slide_height-Inches(0.45), BW, ACCS)

# ============================================================
# MAIN
# ============================================================

STYLE_MAP = {
    "a-clean-professional": render_standard_style,
    "b-dark-cyber": render_b_style,
    "c-warm-business": render_standard_style,
    "d-minimal-ink": render_standard_style,
    "e-bold-vibrant": render_standard_style,
    "f-elegant-classic": render_f_style,
    "g-hardcore-industrial": render_g_style,
    "h-memphis-dopamine": render_h_style,
}

def main():
    if len(sys.argv) < 4:
        print("Usage: python render.py <style_id> <content_md> <output_pptx>")
        print(f"Available styles: {', '.join(STYLE_MAP.keys())}")
        sys.exit(1)

    style_id = sys.argv[1]
    content_file = sys.argv[2]
    output_file = sys.argv[3]

    if style_id not in STYLE_MAP:
        print(f"Unknown style: {style_id}")
        print(f"Available: {', '.join(STYLE_MAP.keys())}")
        sys.exit(1)

    # Load style
    style_path = os.path.join(BASE, "config", "styles", f"{style_id}.yaml")
    if not os.path.exists(style_path):
        print(f"Style file not found: {style_path}")
        sys.exit(1)

    with open(style_path, encoding="utf-8") as f:
        s = yaml.safe_load(f)
    PAL = s["palette"]
    TYPO = s["typography"]

    # Parse content
    blocks = parse_content(content_file)
    print(f"Parsed {len(blocks)} slides from {content_file}")
    print(f"Style: {s['name']} ({style_id})")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Dispatch to style renderer
    render_fn = STYLE_MAP[style_id]
    render_fn(prs, blocks, PAL, TYPO)

    # Save
    os.makedirs(os.path.dirname(output_file) or ".", exist_ok=True)
    prs.save(output_file)
    print(f"Saved: {output_file}")

if __name__ == "__main__":
    main()
